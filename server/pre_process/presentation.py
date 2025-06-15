from pptx import Presentation


def extract_pptx(file_path):
    prs = Presentation(file_path)
    slides = []

    for i, slide in enumerate(prs.slides):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())
        slides.append({"slide_number": i + 1, "content": slide_content})
    return slides
